"""IndexerTextExtractor——从文档附件中提取文本用于全文索引。

对齐 Java IndexerTextExtractor。
支持 PDF、Office、纯文本等格式。
"""
import os
import tempfile
from typing import Dict, Set


class IndexerTextExtractor:
    """文档文本提取器。"""

    # 支持的文件扩展名
    TEXT_EXTENSIONS = {".txt", ".csv", ".log"}
    PDF_EXTENSIONS = {".pdf"}
    WORD_EXTENSIONS = {".doc", ".docx"}
    EXCEL_EXTENSIONS = {".xls", ".xlsx"}
    PPT_EXTENSIONS = {".ppt", ".pptx", ".pps"}
    ODF_EXTENSIONS = {".odt", ".ods", ".odp", ".odg"}

    def extract_texts(self, binary_resources: Set) -> Dict[str, str]:
        """从二进制资源集合中提取文本。

        返回: {full_name: extracted_text}
        """
        result = {}
        for br in binary_resources:
            if not br or not br.full_name:
                continue
            fname = getattr(br, 'name', '') or ''
            ext = os.path.splitext(fname or '')[1].lower()
            try:
                text = self._extract_single(br, ext)
                if text:
                    result[br.full_name] = text
            except Exception:
                pass  # 静默跳过提取失败的文件
        return result

    def _extract_single(self, br, ext: str) -> str:
        """提取单个文件的文本。"""
        data = getattr(br, 'data', None)
        if data is None:
            return ""

        if ext in self.TEXT_EXTENSIONS:
            return self._extract_text(data)

        if ext in self.PDF_EXTENSIONS:
            return self._extract_pdf(data)

        if ext in self.WORD_EXTENSIONS:
            return self._extract_word(data, ext)

        if ext in self.EXCEL_EXTENSIONS:
            return self._extract_excel(data, ext)

        if ext in self.PPT_EXTENSIONS:
            return self._extract_ppt(data, ext)

        return ""

    def _extract_text(self, data: bytes) -> str:
        try:
            return data.decode("utf-8", errors="replace")
        except Exception:
            return data.decode("latin-1", errors="replace")

    def _extract_pdf(self, data: bytes) -> str:
        """使用 PyPDF2 提取 PDF 文本。"""
        try:
            from PyPDF2 import PdfReader
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
                f.write(data)
                f.flush()
                reader = PdfReader(f.name)
                text = ""
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"
            os.unlink(f.name)
            return text
        except ImportError:
            return ""
        except Exception:
            return ""

    def _extract_word(self, data: bytes, ext: str) -> str:
        """提取 Word 文本。"""
        try:
            if ext == ".docx":
                try:
                    from docx import Document
                    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
                        f.write(data)
                        f.flush()
                        doc = Document(f.name)
                        text = "\n".join(p.text for p in doc.paragraphs)
                    os.unlink(f.name)
                    return text
                except ImportError:
                    pass
            return ""
        except Exception:
            return ""

    def _extract_excel(self, data: bytes, ext: str) -> str:
        """提取 Excel 文本。"""
        try:
            if ext == ".xlsx":
                try:
                    from openpyxl import load_workbook
                    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
                        f.write(data)
                        f.flush()
                        wb = load_workbook(f.name, read_only=True)
                        text = ""
                        for ws in wb.worksheets:
                            for row in ws.iter_rows(values_only=True):
                                text += " ".join(str(c) for c in row if c) + "\n"
                    os.unlink(f.name)
                    return text
                except ImportError:
                    pass
            return ""
        except Exception:
            return ""

    def _extract_ppt(self, data: bytes, ext: str) -> str:
        """提取 PPT 文本。"""
        try:
            if ext == ".pptx":
                try:
                    from pptx import Presentation
                    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
                        f.write(data)
                        f.flush()
                        prs = Presentation(f.name)
                        text = ""
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text += shape.text + "\n"
                    os.unlink(f.name)
                    return text
                except ImportError:
                    pass
            return ""
        except Exception:
            return ""


indexer_text_extractor = IndexerTextExtractor()
